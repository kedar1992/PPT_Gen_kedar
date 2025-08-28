from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def analyze_template(template_stream, allow_image_reuse=True):
    prs = Presentation(template_stream)
    layouts = [layout.name for layout in prs.slide_layouts]
    title_font = None
    body_font = None

    try:
        master = prs.slide_masters[0]
        for shape in master.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if run.font.name:
                            if not title_font:
                                title_font = run.font.name
                            body_font = body_font or run.font.name
    except:
        pass

    images = []
    if allow_image_reuse:
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    images.append(shape.image.blob)

    return {
        "layouts": layouts,
        "title_font": title_font,
        "body_font": body_font,
        "images": images[:5]
    }
