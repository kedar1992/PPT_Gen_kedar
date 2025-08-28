import io
from pptx import Presentation
from pptx.util import Pt

def build_presentation(template_stream, plan, style_info, reuse_images=True):
    prs = Presentation(template_stream)
    title_font = style_info.get("title_font")
    body_font = style_info.get("body_font")

    for slide_data in plan["slides"]:
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data["title"]
            if title_font:
                for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                    run.font.name = title_font
                    run.font.size = Pt(36)

        # Bullets
        body = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                body = shape
                break
        if body:
            tf = body.text_frame
            tf.clear()
            for bullet in slide_data.get("bullets", []):
                p = tf.add_paragraph()
                p.text = bullet
                if body_font:
                    for run in p.runs:
                        run.font.name = body_font
                        run.font.size = Pt(18)

        # Speaker notes
        if slide_data.get("notes"):
            notes = slide.notes_slide.notes_text_frame
            notes.text = slide_data["notes"]

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out
